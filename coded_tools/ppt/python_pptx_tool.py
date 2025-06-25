from typing import Dict, Any, Union
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os
from neuro_san.interfaces.coded_tool import CodedTool

class PythonPPTXTool(CodedTool):
    """Implements a tool to create PPTX slides from text."""

    def invoke(self, args: Dict[str, Any], sly_data: Dict[str, Any]) -> Union[Dict[str, Any], str]:
        slides_text = args.get("slides_text", "")
        if not slides_text:
            return "Error: 'slides_text' is required."

        try:
            prs = Presentation()

            # Select layout safely
            layout = next((l for l in prs.slide_layouts if l.name == "Title and Content"), prs.slide_layouts[1])

            # Process slide text
            for slide_text in slides_text.split('$'):
                if not slide_text.strip():
                    continue

                slide = prs.slides.add_slide(layout)
                title, *bullets = slide_text.split('\n')
                slide.shapes.title.text = title.strip()

                content = slide.placeholders[1]
                text_frame = content.text_frame
                text_frame.clear()

                for i, bullet in enumerate(bullets):
                    bullet = bullet.strip()
                    if not bullet:
                        continue

                    p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                    p.text = bullet
                    p.level = 0
                    p.font.size = Pt(15)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0, 0, 0)

            # Output path (user-specified or default)
            output_file = args.get("output_path", r"D:\nihal\neuro-san-studio-main\neuro-san-ppt-output.pptx")
            prs.save(output_file)
            
            print(f"PPTX generated successfully at {output_file}")
            return {"status": "success", "file_path": os.path.abspath(output_file)}

        except Exception as e:
            return f"Error: {str(e)}"

    async def async_invoke(self, args: Dict[str, Any], sly_data: Dict[str, Any]) -> Union[Dict[str, Any], str]:
        return self.invoke(args, sly_data)
